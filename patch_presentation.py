"""
patch_presentation.py

Opens the baseline PPT from OneDrive, updates only the numbers that
changed in the final evaluation run, adds the two required capstone
slides (Objectives Achieved + Lessons Learned), and saves.

Final run numbers:
  Layer 1  : Precision 95.5%  Recall 84.0%  F1 89.4%
  Layer 2  : Precision 100.0% Recall 16.0%  F1 27.6%
  Layer 3  : TP=10  FP=0  TN=10  FN=40  | Precision 100.0%  Recall 20.0%  F1 33.3%
  Ensemble : Precision 95.7%  Recall 90.0%  F1 92.8%
"""

import os, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

SRC  = r"C:\Users\nemur\OneDrive\Pictures\Tendai_Nemure_CS.pptx"
OUT  = r"C:\Users\nemur\Downloads\cot-poisoning-detection\Tendai_Nemure_CS.pptx"

prs = Presentation(SRC)
BLANK = prs.slide_layouts[6]

# ── colours (match baseline) ──────────────────────────────────────────────────
NAVY   = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0xE3, 0x65, 0x14)
GRAY   = RGBColor(0x44, 0x44, 0x44)
LIGHT  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x1B, 0x7A, 0x3E)
RED    = RGBColor(0xC6, 0x28, 0x28)
YELLOW = RGBColor(0xD4, 0x9A, 0x00)
TEAL   = RGBColor(0x00, 0x6D, 0x77)

# ── simple text-replacement helper ───────────────────────────────────────────
def replace_text(slide, old, new):
    """Replace every occurrence of `old` with `new` across all text frames."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)

# ── number patches ────────────────────────────────────────────────────────────
# Slides are 0-indexed. From the baseline:
#   slide[8]  = Layer 3 (slide 9)
#   slide[10] = Evaluation Results (slide 11)
#   slide[11] = Analysis (slide 12)

s_l3   = prs.slides[8]   # Layer 3
s_eval = prs.slides[10]  # Evaluation Results
s_anal = prs.slides[11]  # Analysis

# ── Layer 3 slide ─────────────────────────────────────────────────────────────
# Confusion matrix values
replace_text(s_l3, '"38"', '"10"')  # TP cell value rendered as text
replace_text(s_l3, '"12"', '"40"')  # FN
replace_text(s_l3, '"1"',  '"0"')   # FP
replace_text(s_l3, '"9"',  '"10"')  # TN

# direct string replacements for the cell values
replace_text(s_l3, "38", "10")   # TP
replace_text(s_l3, "12", "40")   # FN  (be careful - only in this slide)

# KPI bar numbers
replace_text(s_l3, "97.4 %", "100.0 %")
replace_text(s_l3, "76.0 %", "20.0 %")
replace_text(s_l3, "85.4 %", "33.3 %")
replace_text(s_l3, "78.3 %", "—")

# insight callout
replace_text(s_l3, "Only 1 false alarm in 60 chains. When it fires — it's almost always right.",
             "0 false alarms in 60 chains. When it fires — it is always right.")
replace_text(s_l3, "Only 1 false alarm", "0 false alarms")

# ── Evaluation Results table ──────────────────────────────────────────────────
# Layer 3 row in the table — these only appear in L3 context on this slide
replace_text(s_eval, "97.4 %", "100.0 %")
replace_text(s_eval, "76.0 %", "20.0 %")
replace_text(s_eval, "85.4 %", "33.3 %")
replace_text(s_eval, "78.3 %", "—")

# ── Analysis slide ────────────────────────────────────────────────────────────
replace_text(s_anal, "Precision 97.4 %", "Precision 100.0 %")
replace_text(s_anal, "Recall 76 %", "Recall 20 %")
replace_text(s_anal, "Only 1 false alarm across 60 chains.", "0 false alarms across 60 chains.")
replace_text(s_anal, "High precision, moderate recall", "High precision, low recall")


# ── helpers for new slides ────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=14, bold=False,
             color=GRAY, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def header_bar(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), NAVY)
    add_text(slide, Inches(0.4), Inches(0.13), Inches(12.5), Inches(0.6),
             title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.35),
                 subtitle, size=13, color=LIGHT, italic=True)
    add_rect(slide, Inches(0), Inches(1.1), prs.slide_width, Inches(0.05), ACCENT)

def footer(slide, page, total=16):
    add_text(slide, Inches(0.3), Inches(7.1), Inches(10.5), Inches(0.3),
             "Detecting Chain-of-Thought Poisoning  |  Tendai Nemure",
             size=10, color=GRAY)
    add_text(slide, Inches(12.2), Inches(7.1), Inches(1.1), Inches(0.3),
             f"{page} / {total}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ── New Slide A — Objectives Achieved ────────────────────────────────────────
def add_objectives_achieved():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Project Outcomes — Objectives Met")

    rows = [
        ("1", "Build a corpus",
         "4,000 Tensor Trust attacks  ·  50 labelled triage cases  ·  250 poisoned test instances"),
        ("2", "Three-layer defence",
         "Layer 1: 95.5 % precision  ·  Layer 2: 100.0 % precision  ·  Layer 3: 100.0 % precision"),
        ("3", "Ensemble risk scoring",
         "Precision 95.7 %  ·  Recall 90.0 %  ·  F1 92.8 %"),
        ("4", "Inline gateway",
         "FastAPI  POST /triage  —  allow / flag / block with per-layer signals"),
        ("5", "Measure everything",
         "Per-layer confusion matrices  ·  Precision / Recall / F1 / Latency  ·  seed=42, temp=0"),
    ]

    for i, (num, title, detail) in enumerate(rows):
        y = Inches(1.38) + Inches(1.0) * i
        add_rect(s, Inches(0.3), y, Inches(0.5), Inches(0.8), ACCENT)
        add_text(s, Inches(0.3), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, Inches(0.85), y, Inches(11.65), Inches(0.8), LIGHT)
        add_text(s, Inches(1.0), y + Inches(0.06), Inches(2.2), Inches(0.38),
                 title, size=14, bold=True, color=NAVY)
        add_text(s, Inches(3.3), y + Inches(0.06), Inches(9.0), Inches(0.68),
                 detail, size=13, color=GRAY)

    add_rect(s, Inches(0.3), Inches(6.72), Inches(12.7), Inches(0.55), NAVY)
    add_text(s, Inches(0.5), Inches(6.79), Inches(12.3), Inches(0.42),
             "Precision 95.7 %  ·  Recall 90.0 %  ·  F1 92.8 %  ·  Latency < 1 ms (Layer 1)  ·  Inline FastAPI gateway deployed",
             size=13, bold=True, color=WHITE)

    footer(s, 15)


# ── New Slide B — Lessons Learned ────────────────────────────────────────────
def add_lessons_learned():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Lessons Learned")

    lessons = [
        ("Poison realism requires calibration", ACCENT,
         "v2 (explicit directives): 72 % drift — unrealistic.   "
         "v3 (hedging only): 0 % drift — model ignored it.   "
         "v3.5 (SOC-native framing): 24.8 % drift — realistic attack surface."),
        ("Ensemble design matters as much as individual layer performance", TEAL,
         "Weighted average scoring → 20 % recall.   "
         "OR logic (any layer fires = flag, 2+ = block) → 90 % recall."),
        ("Model robustness is not sufficient protection alone", NAVY,
         "claude-haiku-4-5 resisted 75 % of poison.   "
         "The 25 % that succeeded are caught by the ensemble.   "
         "Runtime monitoring and model robustness are complements."),
        ("Ground truth methodology is critical", GRAY,
         "Labelling by poison-present is wrong — poison often fails.   "
         "Fix: use Layer 2 behavioral drift as the oracle."),
    ]

    for i, (title, color, detail) in enumerate(lessons):
        y = Inches(1.38) + Inches(1.28) * i
        add_rect(s, Inches(0.3), y, Inches(0.18), Inches(1.1), color)
        add_rect(s, Inches(0.5), y, Inches(12.5), Inches(1.1), LIGHT)
        add_text(s, Inches(0.65), y + Inches(0.07), Inches(12.1), Inches(0.38),
                 title, size=14, bold=True, color=color)
        add_text(s, Inches(0.65), y + Inches(0.48), Inches(12.1), Inches(0.57),
                 detail, size=12, color=GRAY)

    footer(s, 16)


# ── insert new slides before the last slide (Q&A / final slide) ───────────────
# The baseline has 14 slides; we insert at position 13 (before slide 14)
add_objectives_achieved()
add_lessons_learned()

# Move the two new slides to position 13 and 14 (before the last slide)
from pptx.oxml.ns import qn
from lxml import etree

slide_list = prs.slides._sldIdLst
all_ids = list(slide_list)
# last two entries are the slides we just added
new_a = all_ids[-2]
new_b = all_ids[-1]
last  = all_ids[-3]   # was the last slide before we added

# Remove from end, insert before last
slide_list.remove(new_a)
slide_list.remove(new_b)
idx = list(slide_list).index(last)
slide_list.insert(idx, new_b)
slide_list.insert(idx, new_a)

# ── save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
